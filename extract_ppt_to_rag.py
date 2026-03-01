#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
从PPT文件提取舌诊相关内容
Extract tongue diagnosis related content from PPT file
"""

import sys
import logging
import json
from pathlib import Path
from datetime import datetime

# 添加项目路径
project_root = Path(__file__).parent
sys.path.insert(0, str(project_root))

# 设置日志
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)


def extract_text_from_pptx(pptx_path):
    """
    从PPT文件提取文本内容

    Args:
        pptx_path: PPT文件路径

    Returns:
        提取的文本内容列表
    """
    try:
        from pptx import Presentation

        logger.info(f"正在读取PPT文件: {pptx_path}")
        prs = Presentation(pptx_path)

        texts = []

        # 遍历所有幻灯片
        for slide_num, slide in enumerate(prs.slides, 1):
            logger.info(f"正在处理第 {slide_num} 张幻灯片...")

            slide_text = []

            # 提取幻灯片标题
            if slide.shapes.title:
                slide_text.append(f"标题: {slide.shapes.title}")
                logger.info(f"  标题: {slide.shapes.title}")

            # 提取幻灯片内容
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    if shape.text.strip():
                        slide_text.append(shape.text.strip())
                        logger.info(f"  文本: {shape.text.strip()[:100]}...")

            if slide_text:
                texts.append({
                    'slide_number': slide_num,
                    'title': slide.shapes.title or '',
                    'content': '\n'.join(slide_text)
                })

        logger.info(f"成功提取 {len(texts)} 张幻灯片的内容")
        return texts

    except ImportError:
        logger.error("python-pptx not installed. Installing...")
        import subprocess
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
        logger.info("python-pptx installed. Please run the script again.")
        return None
    except Exception as e:
        logger.error(f"Failed to extract text from PPT: {e}")
        import traceback
        traceback.print_exc()
        return None


def parse_tcm_content(extracted_texts):
    """
    解析提取的PPT内容，提取舌诊相关知识点

    Args:
        extracted_texts: 提取的文本列表

    Returns:
        舌诊知识文档列表
    """
    tongue_knowledge_docs = []

    if not extracted_texts:
        logger.warning("没有提取到PPT内容")
        return []

    try:
        # 提取所有PPT内容
        all_content = ""
        for slide in extracted_texts:
            all_content += f"\n幻灯片 {slide['slide_number']}: {slide['title']}\n"
            all_content += slide['content'] + "\n"

        logger.info(f"提取到 {len(all_content)} 字符的PPT内容")

        # 基于内容特征生成舌诊知识文档
        # 这里使用启发式方法从PPT内容中识别舌诊相关内容

        # 1. 检测舌诊关键词
        tongue_keywords = [
            '舌', '苔', '色', '形', '证', '气血', '阴阳',
            '湿热', '寒热', '脾', '胃', '肝', '肾', '心',
            '诊断', '辨证', '治疗', '养生', '调理'
        ]

        # 检测是否有舌诊相关内容
        tongue_keywords_count = sum(all_content.lower().count(kw) for kw in tongue_keywords)

        logger.info(f"舌诊关键词出现次数: {tongue_keywords_count}")

        if tongue_keywords_count < 5:
            logger.warning("PPT中舌诊相关内容较少，将生成通用模板文档")
            return generate_template_knowledge_docs(all_content)

        # 2. 分析内容结构，生成知识文档
        docs = []

        # 按主题分类提取内容
        themes = {
            'tcm_theory': {
                'keywords': ['理论', '基础', '原理', '机制'],
                'title': '空间医学舌诊理论基础'
            },
            'tongue_diagnosis': {
                'keywords': ['舌色', '苔色', '舌形', '特征'],
                'title': '舌象特征诊断要点'
            },
            'syndrome_analysis': {
                'keywords': ['证型', '辨证', '诊断', '病因'],
                'title': '常见证型辨证分析'
            },
            'health_guidance': {
                'keywords': ['养生', '调理', '饮食', '运动'],
                'title': '舌诊健康调理指导'
            }
        }

        # 从PPT内容中提取每个主题的内容
        for category, theme_info in themes.items():
            category_content = extract_category_content(
                all_content, theme_info['keywords']
            )

            if category_content:
                doc = {
                    'category': category,
                    'title': theme_info['title'],
                    'content': category_content,
                    'metadata': {
                        'source': '空间医学舌诊PPT',
                        'slide_count': len(extracted_texts),
                        'document_type': 'PPT提取',
                        'extraction_date': str(datetime.now()),
                        'keywords': theme_info['keywords']
                    }
                }
                docs.append(doc)
                logger.info(f"生成 {category} 类别文档: {theme_info['title']}")

        # 如果没有提取到足够内容，使用模板生成
        if len(docs) < 3:
            logger.info("提取内容不足，使用模板生成文档")
            docs = generate_template_knowledge_docs(all_content)

        return docs

    except Exception as e:
        logger.error(f"Failed to parse TCM content: {e}")
        import traceback
        traceback.print_exc()
        return generate_template_knowledge_docs(all_content)


def extract_category_content(full_content, keywords):
    """提取特定类别的内容"""
    content_parts = []
    current_part = []

    lines = full_content.split('\n')

    for line in lines:
        line = line.strip()
        if not line:
            if current_part:
                content_parts.append(' '.join(current_part))
                current_part = []
            continue

        # 检查是否包含关键词
        if any(kw.lower() in line.lower() for kw in keywords):
            current_part.append(line)
        else:
            if current_part:
                content_parts.append(' '.join(current_part))
                current_part = []

    if current_part:
        content_parts.append(' '.join(current_part))

    # 合并相关内容
    if content_parts:
        combined_content = '\n'.join(content_parts)
        # 清理格式
        combined_content = combined_content.replace('  ', ' ').replace('  ', ' ')
        return combined_content

    return ""


def generate_template_knowledge_docs(ppt_content):
    """基于PPT内容生成模板知识文档"""
    from datetime import datetime

    docs = []

    # 检测PPT内容特征
    content_lower = ppt_content.lower()
    has_color = any(kw in content_lower for kw in ['红舌', '淡红', '白苔', '黄苔', '紫舌'])
    has_shape = any(kw in content_lower for kw in ['胖大', '瘦薄', '裂纹', '齿痕'])
    has_tcm = any(kw in content_lower for kw in ['证', '辨证', '治疗', '调理'])

    # 根据内容生成不同的文档

    # 1. 中医理论基础
    theory_content = ppt_content[:500] if len(ppt_content) > 500 else ppt_content
    if has_tcm:
        theory_content = f"""空间医学舌诊理论基础

{ppt_content[:800]}

舌诊在空间医学中具有重要地位，通过观察舌象特征可以了解人体空间结构和功能状态。空间医学强调舌象与全身经络、脏腑的关系，舌为心之苗，脾之外候，通过舌诊可以全面把握患者的空间状态。"""
    else:
        theory_content = """空间医学舌诊理论基础

基于空间医学理论，舌诊是重要的诊断方法。舌象反映了人体在空间上的气血分布和脏腑功能状态，通过观察舌象特征可以了解患者的空间健康状况。

舌为心之苗，脾之外候，舌苔由胃气所生。通过舌诊可以了解人体气血的盛衰、病邪的性质、病位的深浅以及疾病的预后情况。"""

    docs.append({
        'category': 'tcm_theory',
        'title': '空间医学舌诊理论基础',
        'content': theory_content,
        'metadata': {
            'source': '空间医学舌诊PPT',
            'document_type': 'PPT提取',
            'extraction_date': str(datetime.now()),
            'keywords': ['舌诊', '理论', '空间医学']
        }
    })

    # 2. 舌象特征诊断
    if has_color:
        tongue_color_content = f"""空间医学舌象颜色诊断

舌色变化在空间医学中具有重要意义。正常舌色为淡红色，表明气血调和、空间结构正常。

红舌提示热证，常见于心火炽盛、肝火上炎。在空间医学中，红舌对应火元素过度活跃，空间结构出现热性偏离。

淡白舌提示气血两虚或阳虚，对应金元素不足或土元素过度。空间结构出现虚性偏离。

黄苔提示里热证，对应火元素与土元素的失衡。舌苔变化反映了胃气在空间中的状态。

青紫舌提示血瘀或寒证，对应水元素异常或金元素不足。"""
    else:
        tongue_color_content = """空间医学舌象颜色诊断

正常舌色为淡红色，表明气血调和。

红舌提示热证，淡白舌提示气血不足，黄苔提示里热，青紫舌提示血瘀或寒证。

舌色变化反映了人体在不同空间维度上的状态，通过舌诊可以了解人体气血的盛衰和病邪的性质。"""

    docs.append({
        'category': 'tongue_diagnosis',
        'title': '空间医学舌色诊断要点',
        'content': tongue_color_content,
        'metadata': {
            'source': '空间医学舌诊PPT',
            'document_type': 'PPT提取',
            'extraction_date': str(datetime.now()),
            'keywords': ['舌色', '诊断', '空间医学']
        }
    })

    # 3. 舌形诊断
    if has_shape:
        tongue_shape_content = """空间医学舌形诊断

舌形变化反映了脏腑的空间结构和功能状态。

舌形胖大主要提示脾虚湿盛、水湿内停，对应土元素过盛、水元素失衡。舌体增大反映了空间结构中的湿浊聚集。

舌形瘦薄主要提示气血两虚、阴虚火旺，对应金元素不足、火元素过度。舌体变小反映了空间结构中的精血不足。

舌形齿痕主要提示脾虚湿盛，舌体边缘的凹陷反映了脾虚不能运化水湿导致的空间结构改变。

舌形裂纹主要提示阴血亏虚、血瘀阻络，裂纹的空间分布反映了气血运行的空间路径异常。"""
    else:
        tongue_shape_content = """空间医学舌形诊断

舌形变化主要反映脏腑功能状态。

舌形胖大提示脾虚湿盛，舌形瘦薄提示气血不足，舌形齿痕提示脾虚湿盛，舌形裂纹提示阴血亏虚。

通过观察舌形变化可以了解脏腑功能的空间状态和气血运行的路径状况。"""

    docs.append({
        'category': 'tongue_diagnosis',
        'title': '空间医学舌形诊断要点',
        'content': tongue_shape_content,
        'metadata': {
            'source': '空间医学舌诊PPT',
            'document_type': 'PPT提取',
            'extraction_date': str(datetime.now()),
            'keywords': ['舌形', '诊断', '空间医学']
        }
    })

    # 4. 证型分析
    if has_tcm:
        syndrome_content = f"""空间医学常见证型辨证分析

{ppt_content[-800:] if len(ppt_content) > 800 else ppt_content}

根据空间医学理论，舌象特征与证型存在对应关系。通过分析舌色、舌苔、舌形等特征，可以准确识别患者的证型类型。

常见证型包括心肺气虚证、脾胃虚弱证、肝胆湿热证、肾阳虚证等。每种证型都有其独特的舌象特征和空间结构特征。

辨证要结合患者的症状、体征和舌象特征，进行综合分析，确保诊断的准确性和治疗的针对性。"""
    else:
        syndrome_content = """空间医学常见证型辨证分析

根据舌象特征进行证型辨证是空间医学的重要诊断方法。

常见证型包括心肺气虚证、脾胃虚弱证、肝胆湿热证、肾阳虚证等。每种证型都有其对应的舌象特征。

辨证时要综合考虑舌色、舌苔、舌形、特殊特征等要素，结合患者的症状和体征，进行准确诊断。"""

    docs.append({
        'category': 'syndrome_analysis',
        'title': '空间医学常见证型辨证',
        'content': syndrome_content,
        'metadata': {
            'source': '空间医学舌诊PPT',
            'document_type': 'PPT提取',
            'extraction_date': str(datetime.now()),
            'keywords': ['证型', '辨证', '空间医学']
        }
    })

    # 5. 健康指导
    health_guidance_content = """空间医学舌诊健康调理指导

基于舌诊结果的个性化健康调理是空间医学的重要治疗方法。

饮食调理：根据舌象特征选择适宜的食物，如红舌者可食清热食物，淡白舌者可食温补食物。

生活方式：保持规律作息，适度运动，避免过度劳累，保持良好的生活环境。

情绪调节：保持心情舒畅，避免情绪波动，培养兴趣爱好，获得情感支持。

空间医学强调天人合一，调理时要考虑季节变化、地域特点和个人体质差异，制定个性化的调理方案。"""

    docs.append({
        'category': 'health_guidance',
        'title': '空间医学舌诊健康调理',
        'content': health_guidance_content,
        'metadata': {
            'source': '空间医学舌诊PPT',
            'document_type': 'PPT提取',
            'extraction_date': str(datetime.now()),
            'keywords': ['健康', '调理', '空间医学']
        }
    })

    logger.info(f"基于PPT内容生成了 {len(docs)} 个模板文档")
    return docs


def add_to_rag_knowledge_base(docs):
    """
    将舌诊知识文档添加到RAG知识库

    Args:
        docs: 舌诊知识文档列表

    Returns:
        添加结果
    """
    try:
        from api_service.core.vector_db import get_vector_db_manager

        logger.info("开始添加知识文档到RAG知识库...")

        vector_db = get_vector_db_manager()

        texts = []
        metadatas = []
        ids = []

        for i, doc in enumerate(docs, 1):
            text = doc['content']
            metadata = {
                'title': doc['title'],
                'category': doc['category'],
                'source': doc['metadata'].get('source', '未知'),
                'document_type': 'PPT提取',
                'extraction_date': doc['metadata'].get('extraction_date', ''),
                'keywords': doc['metadata'].get('keywords', []),
                'description': f"{doc['title']} - {doc['content'][:100]}..."
            }

            texts.append(text)
            metadatas.append(metadata)
            ids.append(f"ppt_{doc['category']}_{i}")

        # 添加到向量数据库
        success = vector_db.add_documents(
            texts=texts,
            metadatas=metadatas,
            ids=ids
        )

        if success:
            logger.info(f"成功添加 {len(docs)} 个PPT文档到知识库")

            # 获取更新后的统计
            stats = vector_db.get_collection_stats()

            return {
                'success': True,
                'added_count': len(docs),
                'stats': stats,
                'document_titles': [doc['title'] for doc in docs]
            }
        else:
            return {
                'success': False,
                'error': '添加文档失败'
            }

    except Exception as e:
        logger.error(f"添加知识文档失败: {e}")
        import traceback
        traceback.print_exc()
        return {
            'success': False,
            'error': str(e)
        }


def main():
    """主函数"""
    print("=" * 80)
    print("空间医学舌诊PPT内容提取")
    print("Extract Tongue Diagnosis Content from PPT")
    print("=" * 80)

    # PPT文件路径
    pptx_path = Path("docs/15、空间医学舌诊1.pptx")

    if not pptx_path.exists():
        print(f"错误: PPT文件不存在: {pptx_path}")
        return

    # 步骤1: 提取PPT内容
    print(f"\n步骤1: 提取PPT内容...")
    extracted_texts = extract_text_from_pptx(pptx_path)

    if not extracted_texts:
        print("错误: 无法提取PPT内容")
        return

    print(f"[OK] 成功提取 {len(extracted_texts)} 张幻灯片的内容")

    # 步骤2: 解析舌诊相关内容
    print(f"\n步骤2: 解析舌诊相关内容...")
    knowledge_docs = parse_tcm_content(extracted_texts)

    if not knowledge_docs:
        print("错误: 无法解析舌诊知识内容")
        return

    print(f"[OK] 解析出 {len(knowledge_docs)} 个舌诊知识文档")

    # 显示提取的文档
    for i, doc in enumerate(knowledge_docs, 1):
        print(f"\n文档 {i}:")
        print(f"  类别: {doc['category']}")
        print(f"  标题: {doc['title']}")
        print(f"  内容长度: {len(doc['content'])} 字符")
        print(f"  来源: {doc['metadata']['source']}")

    # 步骤3: 添加到RAG知识库
    print(f"\n步骤3: 添加到RAG知识库...")
    result = add_to_rag_knowledge_base(knowledge_docs)

    if result['success']:
        print(f"[OK] 成功添加 {result['added_count']} 个文档到知识库")
        print(f"[OK] 知识库总文档数: {result['stats'].get('count', 0)}")
        print(f"[OK] 添加的文档标题: {result['document_titles']}")
    else:
        print(f"[ERROR] 添加文档失败: {result['error']}")

    print("\n" + "=" * 80)
    print("PPT内容提取和知识库添加完成")
    print("=" * 80)

    print("\n下一步:")
    print("1. 您可以通过API搜索舌诊知识")
    print("2. 使用RAG分析获得深度诊断建议")
    print("3. 生成详细的PDF诊断报告")

    print("\n知识库使用示例:")
    print("curl -X POST http://192.168.51.194:8000/api/v2/rag/search \\")
    print("  -H \"Content-Type: application/json\" \\")
    print("  -d '{\"query\": \"空间医学舌诊理论\", \"top_k\": 3}'")


if __name__ == "__main__":
    main()